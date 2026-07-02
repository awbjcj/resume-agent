from pathlib import Path

SUPPORTED_SUFFIXES = frozenset({".pdf", ".docx", ".txt", ".md", ".pptx"})


def read_document_text(path: str | Path) -> str:
    """Extract plain text from a supported profile source document."""
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix in (".txt", ".md"):
        return p.read_text(encoding="utf-8")
    if suffix == ".pdf":
        return _read_pdf(p)
    if suffix == ".docx":
        return _read_docx(p)
    if suffix == ".pptx":
        return _read_pptx(p)
    supported = ", ".join(sorted(SUPPORTED_SUFFIXES))
    raise ValueError(f"Unsupported document format: {suffix or '(none)'} (use {supported})")


read_resume_text = read_document_text


def _read_pdf(p: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(p))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _read_docx(p: Path) -> str:
    from docx import Document

    doc = Document(str(p))
    parts = [para.text for para in doc.paragraphs]
    # Many resume templates lay out content in tables; doc.paragraphs skips those.
    parts += [cell.text for table in doc.tables for row in table.rows for cell in row.cells]
    return "\n".join(parts)


def _read_pptx(p: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    slides: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is not None and text_frame.text.strip():
                parts.append(text_frame.text)
        if slide.has_notes_slide:
            notes_text_frame = slide.notes_slide.notes_text_frame
            if notes_text_frame is not None and notes_text_frame.text.strip():
                parts.append(notes_text_frame.text)
        slides.append("\n".join(parts))
    return "\n\n".join(slides)
