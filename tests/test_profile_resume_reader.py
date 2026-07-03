import pytest

from resume_agent.profile.resume_reader import (
    SUPPORTED_SUFFIXES,
    read_document_text,
    read_resume_text,
)


def test_reads_txt(tmp_path):
    f = tmp_path / "resume.txt"
    f.write_text("Ada Lovelace\nEngineer", encoding="utf-8")
    assert read_resume_text(f) == "Ada Lovelace\nEngineer"


def test_reads_docx(tmp_path):
    from docx import Document

    f = tmp_path / "resume.docx"
    doc = Document()
    doc.add_paragraph("Ada Lovelace")
    doc.add_paragraph("Analytical Engines Ltd")
    doc.save(str(f))

    text = read_resume_text(f)
    assert "Ada Lovelace" in text
    assert "Analytical Engines Ltd" in text


def test_unsupported_format_raises(tmp_path):
    f = tmp_path / "resume.rtf"
    f.write_text("nope", encoding="utf-8")
    with pytest.raises(ValueError):
        read_resume_text(f)


def test_read_markdown(tmp_path):
    doc = tmp_path / "notes.md"
    doc.write_text("# Projects\n- Built a compiler", encoding="utf-8")
    assert "Built a compiler" in read_document_text(doc)


def test_read_pptx_slides_and_notes(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    assert title is not None
    title.text = "Migration Case Study"
    notes_text_frame = slide.notes_slide.notes_text_frame
    assert notes_text_frame is not None
    notes_text_frame.text = "Led a team of four"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))

    text = read_document_text(path)
    assert "Migration Case Study" in text
    assert "Led a team of four" in text


def test_read_html(tmp_path):
    doc = tmp_path / "page.html"
    doc.write_text("<h1>Projects</h1><p>Built a compiler</p>", encoding="utf-8")
    assert "Built a compiler" in read_document_text(doc)


def test_read_xlsx(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["Project", "Impact"])
    ws.append(["Pipeline rewrite", "Cut runtime 40%"])
    path = tmp_path / "impact.xlsx"
    wb.save(str(path))

    text = read_document_text(path)
    assert "Pipeline rewrite" in text


def test_supported_suffixes_cover_all_formats():
    assert SUPPORTED_SUFFIXES == frozenset(
        {".pdf", ".docx", ".txt", ".md", ".pptx", ".xlsx", ".html"}
    )
